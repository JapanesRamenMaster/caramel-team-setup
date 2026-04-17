#!/usr/bin/env python3
"""
템플릿 PPTX의 슬라이드를 복제하고 텍스트만 교체하여 새 PPTX 생성.
디자인은 100% 보존.

사용법:
    python3 pptx_create.py --template 템플릿.pptx --spec spec.json --output output.pptx
"""

import argparse
import copy
import json
import os
import sys

from pptx import Presentation
from lxml import etree

R_NS = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
P_NS = 'http://schemas.openxmlformats.org/presentationml/2006/main'
A_NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'


def duplicate_slide(prs, src_index):
    """슬라이드를 복제하여 맨 끝에 추가. 복제된 슬라이드 반환."""
    src_slide = prs.slides[src_index]
    src_layout = src_slide.slide_layout

    # 새 슬라이드 추가 (같은 레이아웃)
    new_slide = prs.slides.add_slide(src_layout)

    # 새 슬라이드의 기존 shape 모두 제거
    spTree = new_slide._element.find(f'{{{P_NS}}}cSld/{{{P_NS}}}spTree')
    for child in list(spTree):
        tag = etree.QName(child).localname
        if tag not in ('nvGrpSpPr', 'grpSpPr'):
            spTree.remove(child)

    # 원본 슬라이드의 shape 복사
    src_spTree = src_slide._element.find(f'{{{P_NS}}}cSld/{{{P_NS}}}spTree')
    for child in src_spTree:
        tag = etree.QName(child).localname
        if tag in ('nvGrpSpPr', 'grpSpPr'):
            continue
        new_child = copy.deepcopy(child)
        spTree.append(new_child)

    # 원본 슬라이드의 배경 복사
    src_bg = src_slide._element.find(f'{{{P_NS}}}cSld/{{{P_NS}}}bg')
    if src_bg is not None:
        new_cSld = new_slide._element.find(f'{{{P_NS}}}cSld')
        old_bg = new_cSld.find(f'{{{P_NS}}}bg')
        if old_bg is not None:
            new_cSld.remove(old_bg)
        new_bg = copy.deepcopy(src_bg)
        new_cSld.insert(0, new_bg)

    # 이미지 관계 복사
    for rel in src_slide.part.rels.values():
        if 'image' in rel.reltype.lower():
            try:
                new_slide.part.relate_to(rel.target_part, rel.reltype)
            except Exception:
                pass

    return new_slide


def replace_text_preserve_format(shape, new_text):
    """shape의 텍스트를 교체하면서 서식 보존."""
    if not shape.has_text_frame:
        return

    tf = shape.text_frame

    # 첫 paragraph, 첫 run의 서식 저장
    saved_rPr = None
    saved_pPr = None

    if tf.paragraphs:
        p = tf.paragraphs[0]
        pPr = p._p.find(f'{{{A_NS}}}pPr')
        if pPr is not None:
            saved_pPr = copy.deepcopy(pPr)
        if p.runs:
            rPr = p.runs[0]._r.find(f'{{{A_NS}}}rPr')
            if rPr is not None:
                saved_rPr = copy.deepcopy(rPr)

    # 텍스트 교체
    tf.clear()

    lines = new_text.split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()

        # paragraph 서식 복원
        if saved_pPr is not None:
            old_pPr = para._p.find(f'{{{A_NS}}}pPr')
            if old_pPr is not None:
                para._p.remove(old_pPr)
            para._p.insert(0, copy.deepcopy(saved_pPr))

        run = para.add_run()
        run.text = line

        # run 서식 복원
        if saved_rPr is not None:
            old_rPr = run._r.find(f'{{{A_NS}}}rPr')
            if old_rPr is not None:
                run._r.remove(old_rPr)
            run._r.insert(0, copy.deepcopy(saved_rPr))


def find_text_shapes(slide):
    """슬라이드의 텍스트 shape을 {이름: shape} 딕셔너리로 반환."""
    result = {}
    for shape in slide.shapes:
        if shape.has_text_frame:
            result[shape.name] = shape
    return result


def build_presentation(template_path, spec, output_path):
    """스펙에 따라 PPTX 생성."""
    prs = Presentation(template_path)
    template_count = len(prs.slides)

    slides_spec = spec.get('slides', [])

    for i, slide_spec in enumerate(slides_spec):
        src_idx = slide_spec['src_slide']
        text_map = slide_spec.get('text_map', {})
        label = slide_spec.get('label', '')

        # 원본 슬라이드 복제
        new_slide = duplicate_slide(prs, src_idx)

        # 텍스트 교체
        text_shapes = find_text_shapes(new_slide)
        for shape_name, new_text in text_map.items():
            if shape_name in text_shapes:
                replace_text_preserve_format(text_shapes[shape_name], new_text)
            else:
                print(f"  [경고] '{shape_name}' 없음 (slide[{src_idx}])", file=sys.stderr)

        # 하이퍼링크 추가
        for shape_name, url in slide_spec.get('hyperlinks', {}).items():
            if shape_name in text_shapes:
                shape = text_shapes[shape_name]
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.text.strip():
                            run.hyperlink.address = url

        # shape 숨기기 (화면 밖으로 이동)
        for shape in new_slide.shapes:
            if shape.name in slide_spec.get('hide_shapes', []):
                shape.left = -9144000

        print(f"  [{i+1}/{len(slides_spec)}] slide[{src_idx}] 복제: {label}", file=sys.stderr)

    # 원본 템플릿 슬라이드 삭제 (뒤에서부터)
    for idx in range(template_count - 1, -1, -1):
        rId = prs.slides._sldIdLst[idx].get(f'{{{R_NS}}}id')
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[idx])

    prs.save(output_path)
    print(f"\nPPTX 생성 완료: {output_path} ({len(slides_spec)}슬라이드)", file=sys.stderr)


def main():
    parser = argparse.ArgumentParser(description='템플릿 슬라이드 복제 기반 PPTX 생성')
    parser.add_argument('--template', required=True, help='템플릿 PPTX 경로')
    parser.add_argument('--spec', required=True, help='슬라이드 스펙 JSON 경로')
    parser.add_argument('--output', '-o', required=True, help='출력 PPTX 경로')
    args = parser.parse_args()

    if not os.path.exists(args.template):
        print(f"오류: 템플릿 없음: {args.template}", file=sys.stderr)
        sys.exit(1)

    with open(args.spec, 'r', encoding='utf-8') as f:
        spec = json.load(f)

    build_presentation(args.template, spec, args.output)


if __name__ == '__main__':
    main()
